# SPDX-License-Identifier: GPL-3.0-or-later
"""PowerPoint text and metadata extraction using python-pptx."""

from __future__ import annotations

import io
import os
import shutil
import subprocess
import tempfile
from collections.abc import Iterator
from contextlib import contextmanager
from datetime import datetime
from pathlib import Path
from typing import Any, BinaryIO

from src.utils.exceptions import ExtractionError

PresentationSource = str | Path | bytes | bytearray | BinaryIO


def _find_office_converter() -> str | None:
    """Return a LibreOffice executable capable of converting legacy PPT files."""
    executable = shutil.which("soffice") or shutil.which("libreoffice")
    if executable:
        return executable

    if os.name == "nt":
        candidates = (
            Path(os.environ.get("PROGRAMFILES", ""))
            / "LibreOffice/program/soffice.exe",
            Path(os.environ.get("PROGRAMFILES(X86)", ""))
            / "LibreOffice/program/soffice.exe",
        )
        for candidate in candidates:
            if candidate.is_file():
                return str(candidate)
    return None


def _convert_legacy_ppt(source_path: Path, output_dir: Path) -> Path:
    """Convert a binary .ppt presentation to .pptx with LibreOffice."""
    converter = _find_office_converter()
    if not converter:
        raise ExtractionError(
            "Legacy PPT files require LibreOffice for conversion. Install "
            "LibreOffice or save the presentation as PPTX, then upload it again."
        )

    try:
        result = subprocess.run(
            [
                converter,
                "--headless",
                "--convert-to",
                "pptx",
                "--outdir",
                str(output_dir),
                str(source_path),
            ],
            capture_output=True,
            text=True,
            timeout=90,
            check=False,
        )
    except (OSError, subprocess.SubprocessError) as exc:
        raise ExtractionError(
            "The legacy PPT file could not be converted. Try saving it as PPTX."
        ) from exc

    converted_path = output_dir / f"{source_path.stem}.pptx"
    if result.returncode != 0 or not converted_path.is_file():
        raise ExtractionError(
            "The legacy PPT file is corrupted or could not be converted to PPTX."
        )
    return converted_path


@contextmanager
def _normalized_source(
    source: PresentationSource,
    extension: str | None = None,
) -> Iterator[Any]:
    """Yield a python-pptx-compatible source, converting legacy PPT if needed."""
    source_extension = extension.lower() if extension else ""
    if not source_extension and isinstance(source, (str, Path)):
        source_extension = Path(source).suffix.lower()

    if source_extension != ".ppt":
        if isinstance(source, (bytes, bytearray)):
            yield io.BytesIO(source)
        else:
            yield source
        return

    with tempfile.TemporaryDirectory(prefix="lecture_ppt_") as temp_dir:
        temp_path = Path(temp_dir)
        if isinstance(source, (str, Path)):
            source_path = Path(source)
        else:
            source_path = temp_path / "presentation.ppt"
            if isinstance(source, (bytes, bytearray)):
                source_path.write_bytes(bytes(source))
            else:
                source.seek(0)
                source_path.write_bytes(source.read())
        yield _convert_legacy_ppt(source_path, temp_path)


def _open_presentation(source: Any) -> Any:
    """Open a presentation and normalize dependency/corruption failures."""
    try:
        from pptx import Presentation
        from pptx.exc import PackageNotFoundError
    except ImportError as exc:
        raise ExtractionError(
            "python-pptx is not installed. Install it with: pip install python-pptx"
        ) from exc

    try:
        return Presentation(source)
    except (PackageNotFoundError, KeyError, ValueError, OSError) as exc:
        raise ExtractionError(
            "This PowerPoint presentation is corrupted or is not a valid PPTX file."
        ) from exc
    except Exception as exc:
        raise ExtractionError(
            "PowerPoint extraction failed. Verify the presentation and try again."
        ) from exc


def _is_bullet(paragraph: Any, shape: Any) -> bool:
    """Detect explicit or placeholder-inherited PowerPoint bullets."""
    properties = paragraph._p.pPr
    if properties is not None:
        bullet_tags = {"buAutoNum", "buBlip", "buChar"}
        no_bullet = False
        for child in properties:
            tag_name = child.tag.rsplit("}", 1)[-1]
            if tag_name in bullet_tags:
                return True
            if tag_name == "buNone":
                no_bullet = True
        if no_bullet:
            return False

    if paragraph.level > 0:
        return True
    if not getattr(shape, "is_placeholder", False):
        return False
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER

        return shape.placeholder_format.type in {
            PP_PLACEHOLDER.BODY,
            PP_PLACEHOLDER.OBJECT,
        }
    except (AttributeError, ValueError):
        return False


def _extract_text_frame(shape: Any) -> list[str]:
    """Extract text paragraphs while retaining bullet hierarchy."""
    lines: list[str] = []
    for paragraph in shape.text_frame.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        if _is_bullet(paragraph, shape):
            lines.append(f"{'  ' * paragraph.level}- {text}")
        else:
            lines.append(text)
    return lines


def _extract_table(shape: Any) -> list[str]:
    """Extract a table as pipe-delimited rows."""
    rows: list[str] = []
    for row in shape.table.rows:
        values = [cell.text.strip().replace("\n", " ") for cell in row.cells]
        rows.append(" | ".join(values))
    return rows


def _extract_shape(shape: Any) -> list[str]:
    """Extract supported content from a shape, including grouped shapes."""
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        return []

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        lines: list[str] = []
        for child_shape in shape.shapes:
            lines.extend(_extract_shape(child_shape))
        return lines
    if getattr(shape, "has_table", False):
        return _extract_table(shape)
    if getattr(shape, "has_text_frame", False):
        return _extract_text_frame(shape)
    return []


def _extract_speaker_notes(slide: Any) -> list[str]:
    """Extract speaker notes when the presentation contains them."""
    try:
        if not slide.has_notes_slide:
            return []
        text_frame = slide.notes_slide.notes_text_frame
        return [
            paragraph.text.strip()
            for paragraph in text_frame.paragraphs
            if paragraph.text.strip()
        ]
    except (AttributeError, KeyError, ValueError):
        return []


def _format_datetime(value: datetime | None) -> str:
    return value.isoformat() if value else ""


def _extract_presentation(source: Any) -> tuple[str, dict[str, Any]]:
    """Extract ordered text and metadata from an opened presentation source."""
    presentation = _open_presentation(source)
    properties = presentation.core_properties
    slide_sections: list[str] = []
    preview_parts: list[str] = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_lines = [f"Slide {slide_number}"]
        title_shape = slide.shapes.title
        if title_shape is not None and title_shape.text.strip():
            title = title_shape.text.strip()
            slide_lines.append(f"Title: {title}")
            preview_parts.append(title)

        for shape in slide.shapes:
            if title_shape is not None and shape == title_shape:
                continue
            extracted_lines = _extract_shape(shape)
            slide_lines.extend(extracted_lines)
            preview_parts.extend(line.lstrip("- ") for line in extracted_lines)

        notes = _extract_speaker_notes(slide)
        if notes:
            slide_lines.append("Speaker Notes:")
            slide_lines.extend(notes)
            preview_parts.extend(notes)
        slide_sections.append("\n".join(slide_lines))

    title = (properties.title or "").strip()
    author = (properties.author or "").strip()
    created = _format_datetime(properties.created)
    body_text = "\n\n".join(slide_sections).strip()
    preview = " ".join(preview_parts).strip()[:300]
    metadata = {
        "slide_count": len(presentation.slides),
        "title": title,
        "author": author,
        "created": created,
        "preview": preview,
    }
    metadata_lines = [
        "Presentation Metadata",
        f"Title: {title or 'Unknown'}",
        f"Author: {author or 'Unknown'}",
        f"Created: {created or 'Unknown'}",
        f"Slides: {metadata['slide_count']}",
    ]
    return "\n".join(metadata_lines + ["", body_text]).strip(), metadata


def extract_presentation(
    source: PresentationSource,
    extension: str | None = None,
) -> tuple[str, dict[str, Any]]:
    """Return presentation text and metadata from PPTX or converted PPT input."""
    with _normalized_source(source, extension) as normalized_source:
        return _extract_presentation(normalized_source)


def extract_text_from_pptx(file_path: str) -> str:
    """Extract ordered plain text from a PPTX or legacy PPT presentation."""
    text, _ = extract_presentation(file_path)
    return text


def extract_presentation_metadata(
    source: PresentationSource,
    extension: str | None = None,
) -> dict[str, Any]:
    """Extract presentation metadata for validation and upload preview."""
    _, metadata = extract_presentation(source, extension)
    return metadata
